from typing import Dict,Any
import os
import json
import fitz
import pytesseract
from PIL import Image
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from typing import Union, Dict, Any
from transcribe import WhisperTranscriber
from utils import *
def get_youtube_transcript(youtube_url: str) -> str:
    try:
        existing_subtitles = extract_subtitles(youtube_url)
        if existing_subtitles:
            print("Subtitles found ... Extracted Subtitles from YouTube Video...")
            return existing_subtitles
        else:
            print("Getting YouTube Video's Audio...")
            process_youtube_video(youtube_url)
            audio_file = r"Saved_Media\audio.mp3"
            if not os.path.exists(audio_file):
                raise FileNotFoundError(f"Audio file not found: {audio_file}")
            
            print("Transcribing Audio...")
            transcriber = WhisperTranscriber()
            transcribed_text = transcriber.transcribe(audio_file)
            return transcribed_text
    except Exception as e:
        print(f"Error: {e}")
        return ""


class DocumentExtractor:
    """A unified class for extracting content from various document types including web content."""
    
    def __init__(self):
        """Initialize the document extractor with necessary components."""
        self.youtube_transcriber = WhisperTranscriber()
        
    def extract_content(self, source: str) -> Dict[str, Any]:
        """
        Extract content from various document types including web links.
        
        Args:
            source (str): File path or URL to the content
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - 'content': The extracted text/content
                - 'metadata': Additional information about the extraction
                - 'type': The type of content that was processed
        """
        try:
            # Check if source is a URL
            if source.startswith(('http://', 'https://')):
                if 'youtube.com' in source or 'youtu.be' in source:
                    return self._extract_youtube_content(source)
                elif 'medium.com' in source:
                    return self._extract_medium_content(source)
                else:
                    raise ValueError("Unsupported URL type")
            
            # Handle local files
            if not os.path.exists(source):
                raise FileNotFoundError(f"The file {source} does not exist.")
                
            file_extension = os.path.splitext(source)[1].lower()
            
            handlers = {
                '.txt': self._extract_text_file,
                '.pdf': self._extract_pdf,
                '.pptx': self._extract_pptx
            }
            
            handler = handlers.get(file_extension)
            if not handler:
                raise ValueError(f"Unsupported file type: {file_extension}")
                
            return handler(source)
            
        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': str(e)},
                'type': 'error'
            }
    
    def _extract_text_file(self, file_path: str) -> Dict[str, Any]:
        """Extract content from text files."""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return {
            'content': content,
            'metadata': {'file_size': os.path.getsize(file_path)},
            'type': 'text'
        }
    
    def _extract_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PDF files with OCR support."""
        doc = fitz.open(file_path)
        text_contents = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")
            
            if not text.strip():
                # Try OCR if no text is found
                zoom = 2.0
                matrix = fitz.Matrix(zoom, zoom)
                pixmap = page.get_pixmap(matrix=matrix)
                image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                text = pytesseract.image_to_string(image)
                
            text_contents.append(text)
            
        return {
            'content': '\n'.join(text_contents),
            'metadata': {'pages': len(doc)},
            'type': 'pdf'
        }
    
    def _extract_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PowerPoint presentations."""
        presentation = Presentation(file_path)
        text_contents = []
        
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_contents.append(' '.join(slide_text))
            
        return {
            'content': '\n'.join(text_contents),
            'metadata': {'slides': len(presentation.slides)},
            'type': 'pptx'
        }
    
    def _extract_youtube_content(self, url: str) -> Dict[str, Any]:
        """Extract transcript from YouTube videos."""
        try:
            transcript = get_youtube_transcript(url)
            return {
                'content': transcript,
                'metadata': {'url': url},
                'type': 'youtube'
            }
        except Exception as e:
            raise RuntimeError(f"Failed to extract YouTube content: {str(e)}")
    
    def _extract_medium_content(self, url: str) -> Dict[str, Any]:
        """Extract content from Medium articles."""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Extract title
            title = soup.find('h1').get_text() if soup.find('h1') else ''
            
            # Extract article content
            article_content = []
            content_elements = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
            
            for element in content_elements:
                text = element.get_text().strip()
                if text:
                    article_content.append(text)
            
            return {
                'content': '\n'.join(article_content),
                'metadata': {
                    'title': title,
                    'url': url
                },
                'type': 'medium'
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to extract Medium content: {str(e)}")
